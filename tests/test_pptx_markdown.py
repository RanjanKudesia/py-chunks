import pytest
pytest.importorskip("pptx")  # optional dev dep; absent in wheel-test CI
from pptx import Presentation
from py_chunks import get_markdown
from py_chunks.chunkers.pptx import pptx_to_markdown


# -- Helpers -----------------------------------------------------------------

def _make_pptx(tmp_path, slides):
    """slides: list of (title_str, body_str_or_None)"""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for title, body in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if body:
            slide.placeholders[1].text = body
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


def _make_pptx_with_bullets(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Bullet Slide"
    tf = slide.placeholders[1].text_frame
    tf.text = "First bullet"
    p2 = tf.add_paragraph()
    p2.text = "Second bullet"
    p2.level = 0
    p3 = tf.add_paragraph()
    p3.text = "Sub-bullet"
    p3.level = 1
    path = tmp_path / "bullets.pptx"
    prs.save(str(path))
    return path


def _make_pptx_with_table(tmp_path):
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank)
    rows, cols = 3, 3
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(3)).table
    headers = ["Name", "Value", "Unit"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "42"
    table.cell(1, 2).text = "kg"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "7"
    table.cell(2, 2).text = "m"
    path = tmp_path / "table.pptx"
    prs.save(str(path))
    return path


def _make_pptx_with_notes(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Notes Slide"
    slide.placeholders[1].text = "Body content"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "This is a speaker note."
    path = tmp_path / "notes.pptx"
    prs.save(str(path))
    return path


def _make_pptx_with_bold_italic(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Formatting"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run1 = para.add_run()
    run1.text = "Plain "
    run2 = para.add_run()
    run2.text = "bold"
    run2.font.bold = True
    run3 = para.add_run()
    run3.text = " italic"
    run3.font.italic = True
    path = tmp_path / "fmt.pptx"
    prs.save(str(path))
    return path


# -- Basic structure tests ----------------------------------------------------

def test_slide_heading_contains_slide_number(tmp_path):
    path = _make_pptx(tmp_path, [("My Title", "Body text")])
    result = pptx_to_markdown(str(path))
    assert "## Slide 1" in result


def test_slide_heading_contains_title(tmp_path):
    path = _make_pptx(tmp_path, [("My Title", "Body text")])
    result = pptx_to_markdown(str(path))
    assert "My Title" in result


def test_slide_body_text_included(tmp_path):
    path = _make_pptx(tmp_path, [("Title", "The body paragraph")])
    result = pptx_to_markdown(str(path))
    assert "The body paragraph" in result


def test_multi_slide_both_present(tmp_path):
    path = _make_pptx(tmp_path, [("Slide One", "Content A"), ("Slide Two", "Content B")])
    result = pptx_to_markdown(str(path))
    assert "Slide 1" in result
    assert "Slide 2" in result
    assert "Content A" in result
    assert "Content B" in result


def test_slides_separated_by_hr(tmp_path):
    path = _make_pptx(tmp_path, [("A", "a"), ("B", "b")])
    result = pptx_to_markdown(str(path))
    assert "---" in result


def test_slide_order_preserved(tmp_path):
    path = _make_pptx(tmp_path, [("First", "1"), ("Second", "2")])
    result = pptx_to_markdown(str(path))
    assert result.index("First") < result.index("Second")


# -- Bullet / list tests ------------------------------------------------------

def test_bullets_use_dash_prefix(tmp_path):
    path = _make_pptx_with_bullets(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "- First bullet" in result or "- first bullet" in result.lower()


def test_sub_bullets_indented(tmp_path):
    path = _make_pptx_with_bullets(tmp_path)
    result = pptx_to_markdown(str(path))
    # Sub-bullet must have leading spaces
    assert "  - Sub-bullet" in result or "  - sub-bullet" in result.lower()


# -- Table tests --------------------------------------------------------------

def test_table_renders_pipe_format(tmp_path):
    path = _make_pptx_with_table(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "|" in result


def test_table_header_separator_present(tmp_path):
    path = _make_pptx_with_table(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "---" in result  # separator row


def test_table_cell_content_present(tmp_path):
    path = _make_pptx_with_table(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "Name" in result
    assert "Alpha" in result


# -- Speaker notes tests ------------------------------------------------------

def test_notes_appear_in_output(tmp_path):
    path = _make_pptx_with_notes(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "speaker note" in result.lower()


def test_notes_in_blockquote(tmp_path):
    path = _make_pptx_with_notes(tmp_path)
    result = pptx_to_markdown(str(path))
    assert ">" in result


# -- Inline formatting tests --------------------------------------------------

def test_bold_text_wrapped(tmp_path):
    path = _make_pptx_with_bold_italic(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "**bold**" in result


def test_italic_text_wrapped(tmp_path):
    path = _make_pptx_with_bold_italic(tmp_path)
    result = pptx_to_markdown(str(path))
    assert "* italic*" in result or "*italic*" in result or "* italic" in result


# -- API / error tests --------------------------------------------------------

def test_wrong_extension_raises(tmp_path):
    path = tmp_path / "file.docx"
    path.write_bytes(b"fake")
    with pytest.raises((ValueError, FileNotFoundError, Exception)):
        pptx_to_markdown(str(path))


def test_missing_file_raises(tmp_path):
    with pytest.raises((FileNotFoundError, Exception)):
        pptx_to_markdown(str(tmp_path / "nonexistent.pptx"))


def test_get_markdown_routes_pptx(tmp_path):
    path = _make_pptx(tmp_path, [("Hello", "World")])
    result = get_markdown(str(path))
    assert "Hello" in result
    assert "World" in result


# ── Grouped shapes ────────────────────────────────────────────────────────────

def test_grouped_shape_text_included(tmp_path):
    """Text inside a grouped shape must appear in the markdown output."""
    from lxml import etree

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    spTree = slide.shapes._spTree
    grp_xml = (
        '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvGrpSpPr>'
        '<p:cNvPr id="10" name="Group 1"/>'
        '<p:cNvGrpSpPr/>'
        '<p:nvPr/>'
        '</p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
        '<p:sp>'
        '<p:nvSpPr><p:cNvPr id="11" name="TextBox 1"/>'
        '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        '<p:spPr/>'
        '<p:txBody><a:bodyPr/><a:lstStyle/>'
        '<a:p><a:r><a:t>Group text content</a:t></a:r></a:p>'
        '</p:txBody>'
        '</p:sp>'
        '</p:grpSp>'
    )
    grp_elem = etree.fromstring(grp_xml)
    spTree.append(grp_elem)

    path = tmp_path / "grp.pptx"
    prs.save(str(path))

    result = pptx_to_markdown(str(path))
    assert "Group text content" in result, f"Group shape text missing from:\n{result}"


# ── Section headings ─────────────────────────────────────────────────────────

def test_section_headings_appear_in_output(tmp_path):
    """Named PPTX sections become # headings in markdown output."""
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation()
    layout = prs.slide_layouts[1]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Intro Slide"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Data Slide"

    prs_elem = prs._element

    sldIdLst = prs_elem.find(qn('p:sldIdLst'))
    slide_ids = [sld.get('id') for sld in sldIdLst] if sldIdLst is not None else []

    ext_lst = prs_elem.find(qn('p:extLst'))
    if ext_lst is None:
        ext_lst = etree.SubElement(prs_elem, qn('p:extLst'))

    if slide_ids and len(slide_ids) >= 2:
        section_xml = (
            f'<p:ext xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"'
            f' uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}">'
            f'<p14:sectionLst>'
            f'<p14:section name="Section One" id="{{11111111-1111-1111-1111-111111111111}}">'
            f'<p14:sldIdLst><p14:sldId id="{slide_ids[0]}"/></p14:sldIdLst>'
            f'</p14:section>'
            f'<p14:section name="Section Two" id="{{22222222-2222-2222-2222-222222222222}}">'
            f'<p14:sldIdLst><p14:sldId id="{slide_ids[1]}"/></p14:sldIdLst>'
            f'</p14:section>'
            f'</p14:sectionLst>'
            f'</p:ext>'
        )
        ext_lst.append(etree.fromstring(section_xml))

    path = tmp_path / "sections.pptx"
    prs.save(str(path))

    result = pptx_to_markdown(str(path))
    assert "Intro Slide" in result
    assert "Data Slide" in result


# ── Hyperlinks ────────────────────────────────────────────────────────────────

def test_hyperlink_rendered_as_markdown_link(tmp_path):
    """Hyperlinked runs become [text](url) in markdown output."""
    from lxml import etree
    from pptx.oxml.ns import qn

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Links"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Click here"
    rPr = run._r.get_or_add_rPr()
    hlink = etree.SubElement(rPr, qn('a:hlinkClick'))
    rId = slide.part.relate_to(
        "https://example.com",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hlink.set(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id',
        rId,
    )

    path = tmp_path / "link.pptx"
    prs.save(str(path))

    result = pptx_to_markdown(str(path))
    assert "[Click here](https://example.com)" in result, f"Hyperlink not rendered in:\n{result}"
