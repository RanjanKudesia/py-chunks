"""
Comprehensive quality and performance tests for PPTX chunking (all 7 modes).

Test files (test_files/pptx/):
  minimal.pptx       — 3 slides: title + 2 content slides
  sectioned.pptx     — 11 slides with title-divider section heuristic
  mixed_content.pptx — 8 slides: tables, bullets, prose, short text
  prose_heavy.pptx   — 10 slides: long prose (semantic merge target)
  large_deck.pptx    — 30 slides for performance benchmarking
  empty_slides.pptx  — 6 slides (2 blank) for empty-slide handling
  sample_pptx_generated.pptx — pre-existing fixture

Mode coverage:
  default / structural / semantic / section /
  sliding_window / sentence / page_aware
"""
# pylint: disable=missing-function-docstring,missing-class-docstring
# pylint: disable=redefined-outer-name

import time
from io import BytesIO
from pathlib import Path

import pytest

from py_chunks import (  # pylint: disable=no-name-in-module
    get_chunks,
    get_chunks_from_bytes,
    get_chunks_from_fileobj,
    get_chunks_from_upload,
    stream_chunks,
)
from py_chunks.chunkers.pptx import chunk_pptx, stream_chunk_pptx  # pylint: disable=no-name-in-module


class _DummyUpload:
    def __init__(self, data: bytes, filename: str):
        self.filename = filename
        self.file = BytesIO(data)

# ── Fixture discovery ─────────────────────────────────────────────────────────

_PPTX_DIR = Path(__file__).resolve().parents[2] / "test_files" / "pptx"
ALL_PPTX = sorted(_PPTX_DIR.glob("*.pptx"))

if not ALL_PPTX:
    pytest.skip("No PPTX test files found in test_files/pptx/", allow_module_level=True)

PPTX_IDS = [f.name for f in ALL_PPTX]

# Subsets for parametrised classes
NON_MINIMAL = [f for f in ALL_PPTX if f.name != "minimal.pptx"]
MULTI_SLIDE  = [f for f in ALL_PPTX if f.name in (
    "sectioned.pptx", "mixed_content.pptx", "prose_heavy.pptx",
    "large_deck.pptx", "sectioned.pptx",
)]

STANDARD_KEYS = {"content", "content_type", "metadata"}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _get(path: Path, mode: str = "default", **kw) -> list[dict]:
    return get_chunks(str(path), mode=mode, **kw)


def _timed_get(path: Path, mode: str, **kw):
    """Returns (chunks, rust_ms, python_ms, wall_ms)."""
    t0 = time.perf_counter()
    chunks, timing = chunk_pptx(str(path), mode=mode, **kw)
    wall_ms = round((time.perf_counter() - t0) * 1000, 3)
    return chunks, timing["rust_ms"], timing["python_ms"], wall_ms


def _assert_schema(chunks, label=""):
    tag = f" [{label}]" if label else ""
    assert len(chunks) > 0, f"Expected ≥1 chunk{tag}"
    for c in chunks:
        assert STANDARD_KEYS.issubset(c.keys()), \
            f"Missing schema keys{tag}: {STANDARD_KEYS - c.keys()}"
        assert isinstance(c["content"], str) and c["content"].strip(), \
            f"content must be non-empty string{tag}"
        assert isinstance(c["content_type"], str) and c["content_type"], \
            f"content_type must be non-empty string{tag}"
        assert isinstance(c["metadata"], dict), \
            f"metadata must be dict{tag}"


def _no_duplicate_content(chunks, label=""):
    """All chunk contents should be unique (no verbatim duplication)."""
    seen = set()
    for c in chunks:
        assert c["content"] not in seen, \
            f"Duplicate content in {label}: {c['content'][:60]!r}"
        seen.add(c["content"])


def _avg_chunk_len(chunks) -> float:
    return sum(len(c["content"]) for c in chunks) / max(len(chunks), 1)


# ═══════════════════════════════════════════════════════════════════════════════
# 1. SCHEMA / OUTPUT CONTRACT — every mode on every file
# ═══════════════════════════════════════════════════════════════════════════════

class TestSchemaAllModes:
    """Schema contract: content / content_type / metadata present on every chunk."""

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_default(self, path):
        chunks = _get(path, "default")
        _assert_schema(chunks, f"default/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_structural(self, path):
        chunks = _get(path, "structural")
        _assert_schema(chunks, f"structural/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_semantic(self, path):
        chunks = _get(path, "semantic")
        _assert_schema(chunks, f"semantic/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_section(self, path):
        chunks = _get(path, "section")
        _assert_schema(chunks, f"section/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sliding_window(self, path):
        chunks = _get(path, "sliding_window", window_size=3, overlap=1)
        _assert_schema(chunks, f"sliding_window/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sentence(self, path):
        chunks = _get(path, "sentence", sentences_per_chunk=3)
        _assert_schema(chunks, f"sentence/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_page_aware(self, path):
        chunks = _get(path, "page_aware", paragraphs_per_page=3)
        _assert_schema(chunks, f"page_aware/{path.name}")


# ═══════════════════════════════════════════════════════════════════════════════
# 2. CONTENT_TYPE VALUES
# ═══════════════════════════════════════════════════════════════════════════════

VALID_CONTENT_TYPES = {
    "plain_paragraph", "heading", "bullet_list", "table",
    "long_single_paragraph", "short_disconnected_paragraph",
    "semantic", "section", "sliding_window", "sentence", "page_aware",
}

class TestContentTypes:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_default_content_types_are_valid(self, path):
        for c in _get(path, "default"):
            assert c["content_type"] in VALID_CONTENT_TYPES, \
                f"Unknown content_type '{c['content_type']}' in {path.name}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_semantic_chunks_have_semantic_type(self, path):
        for c in _get(path, "semantic"):
            assert c["content_type"] == "semantic"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sliding_window_type(self, path):
        for c in _get(path, "sliding_window", window_size=2, overlap=1):
            assert c["content_type"] == "sliding_window"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sentence_type(self, path):
        for c in _get(path, "sentence", sentences_per_chunk=2):
            assert c["content_type"] == "sentence"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_page_aware_type(self, path):
        for c in _get(path, "page_aware", paragraphs_per_page=2):
            assert c["content_type"] == "page_aware"


# ═══════════════════════════════════════════════════════════════════════════════
# 3. METADATA CONTRACT PER MODE
# ═══════════════════════════════════════════════════════════════════════════════

class TestMetadataDefault:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_number_present_and_positive(self, path):
        for c in _get(path, "default"):
            m = c["metadata"]
            assert "slide_number" in m, "Missing slide_number"
            assert isinstance(m["slide_number"], int) and m["slide_number"] >= 1

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_range_present_and_consistent_with_slide_number(self, path):
        """slide_range should be present and equal [slide_number, slide_number]."""
        for c in _get(path, "default"):
            m = c["metadata"]
            assert "slide_range" in m, "Missing slide_range in structural/default metadata"
            r = m["slide_range"]
            sn = m["slide_number"]
            assert r == [sn, sn], f"Expected slide_range=[{sn},{sn}], got {r}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_total_slides_consistent(self, path):
        chunks = _get(path, "default")
        totals = {c["metadata"]["document_metadata"]["total_slides"] for c in chunks}
        assert len(totals) == 1, "total_slides not consistent across chunks"
        total = totals.pop()
        assert isinstance(total, int) and total > 0

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_numbers_in_range(self, path):
        chunks = _get(path, "default")
        total = chunks[0]["metadata"]["document_metadata"]["total_slides"]
        for c in chunks:
            sn = c["metadata"]["slide_number"]
            assert 1 <= sn <= total, f"slide_number {sn} out of range [1,{total}]"


class TestMetadataSemantic:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_required_keys(self, path):
        required = {"slide_range", "slide_count", "chunk_index", "document_metadata"}
        for c in _get(path, "semantic"):
            assert required.issubset(c["metadata"].keys()), \
                f"Missing semantic metadata keys: {required - c['metadata'].keys()}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_range_ordering(self, path):
        for c in _get(path, "semantic"):
            r = c["metadata"]["slide_range"]
            assert r[0] <= r[1], f"slide_range start > end: {r}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_chunk_index_sequential(self, path):
        chunks = _get(path, "semantic")
        for i, c in enumerate(chunks):
            assert c["metadata"]["chunk_index"] == i, \
                f"chunk_index {c['metadata']['chunk_index']} != expected {i}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_keyword_density_non_negative(self, path):
        for c in _get(path, "semantic"):
            kd = c["metadata"].get("keyword_density", 0)
            assert kd >= 0, "keyword_density must be ≥ 0"


class TestMetadataSection:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_section_heading_in_metadata(self, path):
        for c in _get(path, "section"):
            m = c["metadata"]
            assert "section_heading" in m, "Missing section_heading in section metadata"
            assert isinstance(m["section_heading"], str) and m["section_heading"]

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_range_present(self, path):
        for c in _get(path, "section"):
            r = c["metadata"].get("slide_range")
            assert r is not None, "Missing slide_range in section chunk"
            assert r[0] <= r[1], f"slide_range inverted: {r}"


class TestMetadataSlidingWindow:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_window_metadata(self, path):
        ws, ov = 3, 1
        for c in _get(path, "sliding_window", window_size=ws, overlap=ov):
            m = c["metadata"]
            assert m["window_size"] == ws
            assert m["overlap"] == ov
            assert "slide_range" in m
            assert "window_index" in m

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_window_index_sequential(self, path):
        chunks = _get(path, "sliding_window", window_size=2, overlap=1)
        for i, c in enumerate(chunks):
            assert c["metadata"]["window_index"] == i


class TestMetadataSentence:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sentence_metadata_keys(self, path):
        required = {"sentences_per_chunk", "actual_sentence_count", "chunk_index", "source_slide"}
        for c in _get(path, "sentence", sentences_per_chunk=3):
            assert required.issubset(c["metadata"].keys())

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_actual_sentence_count_lte_requested(self, path):
        n = 3
        for c in _get(path, "sentence", sentences_per_chunk=n):
            actual = c["metadata"]["actual_sentence_count"]
            assert actual <= n, f"actual_sentence_count {actual} > requested {n}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_sentence_chunk_index_sequential(self, path):
        chunks = _get(path, "sentence", sentences_per_chunk=2)
        for i, c in enumerate(chunks):
            assert c["metadata"]["chunk_index"] == i


class TestMetadataPageAware:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_page_aware_metadata_keys(self, path):
        required = {"slides_per_chunk", "slide_count", "slide_range", "chunk_index"}
        for c in _get(path, "page_aware", paragraphs_per_page=3):
            assert required.issubset(c["metadata"].keys())

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_count_lte_slides_per_chunk(self, path):
        n = 3
        for c in _get(path, "page_aware", paragraphs_per_page=n):
            sc = c["metadata"]["slide_count"]
            assert sc <= n, f"slide_count {sc} > slides_per_chunk {n}"


# ═══════════════════════════════════════════════════════════════════════════════
# 4. CONTENT QUALITY
# ═══════════════════════════════════════════════════════════════════════════════

class TestContentQuality:

    # --- No duplicate chunks ---

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_no_duplicate_content_default(self, path):
        _no_duplicate_content(_get(path, "default"), f"default/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_no_duplicate_content_semantic(self, path):
        _no_duplicate_content(_get(path, "semantic"), f"semantic/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_no_duplicate_content_section(self, path):
        _no_duplicate_content(_get(path, "section"), f"section/{path.name}")

    # --- Chunk length bounds ---

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_default_chunks_not_excessively_long(self, path):
        max_chars = 2000  # structural should stay under ~1200 + MIN merge overhead
        for c in _get(path, "default"):
            assert len(c["content"]) <= max_chars, \
                f"Chunk too long ({len(c['content'])} chars) in {path.name}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_semantic_chunks_under_limit(self, path):
        max_chars = 1600  # Rust cap is 1500 chars plus newline overhead
        for c in _get(path, "semantic"):
            assert len(c["content"]) <= max_chars, \
                f"Semantic chunk too long ({len(c['content'])} chars)"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_page_aware_chunks_under_limit(self, path):
        # paragraphs_per_page=2 ⇒ at most 2 slides worth of text
        chunks = _get(path, "page_aware", paragraphs_per_page=2)
        for c in chunks:
            # Very generous upper bound — just ensure no runaway merges
            assert len(c["content"]) <= 5000, \
                f"page_aware chunk suspiciously large ({len(c['content'])} chars)"

    # --- Table detection ---

    def test_table_slide_gets_table_content_type(self):
        path = _PPTX_DIR / "mixed_content.pptx"
        if not path.exists():
            pytest.skip("mixed_content.pptx not found")
        chunks = _get(path, "default")
        table_types = {c["content_type"] for c in chunks}
        assert "table" in table_types, "Expected at least one 'table' content_type"

    def test_table_chunk_contains_pipe_separated_cells(self):
        """After table-text extraction, cell content should be pipe-separated."""
        path = _PPTX_DIR / "mixed_content.pptx"
        if not path.exists():
            pytest.skip("mixed_content.pptx not found")
        chunks = _get(path, "default")
        table_chunks = [c for c in chunks if c["content_type"] == "table"]
        assert table_chunks, "Expected at least one table chunk"
        assert any("|" in c["content"] for c in table_chunks), (
            "Table chunks should contain pipe-separated cell text"
        )

    # --- Sliding window truncation flag ---

    def test_sliding_window_truncated_flag_is_bool(self):
        """Every sliding-window chunk must carry a 'truncated' boolean metadata field."""
        path = ALL_PPTX[0]
        for c in _get(path, "sliding_window", window_size=2, overlap=0):
            assert "truncated" in c["metadata"], (
                "Missing 'truncated' key in sliding_window metadata"
            )
            assert isinstance(c["metadata"]["truncated"], bool)

    def test_sliding_window_not_truncated_for_normal_decks(self):
        """Normal-sized decks should not trigger the content safety cap."""
        path = ALL_PPTX[0]
        for c in _get(path, "sliding_window", window_size=3, overlap=1):
            assert c["metadata"]["truncated"] is False, (
                "Small deck chunks should not be truncated"
            )

    # --- Speaker notes ---

    def test_notes_text_included_when_present(self):
        """Slides with speaker notes should include [Notes] section in content."""
        import io  # pylint: disable=import-outside-toplevel
        from pptx import Presentation  # pylint: disable=import-outside-toplevel,import-error

        # Build a minimal PPTX with speaker notes in memory.
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Slide"
        slide.shapes.placeholders[1].text = "Body text"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "This is a speaker note."

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        import tempfile  # pylint: disable=import-outside-toplevel
        import os  # pylint: disable=import-outside-toplevel
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(buf.read())
            tmp_path = tmp.name
        try:
            chunks = _get(tmp_path, "default")
            # At least one chunk should contain the notes separator.
            combined = " ".join(c["content"] for c in chunks)
            assert "[Notes]" in combined, "Speaker notes should be included with [Notes] marker"
            assert "speaker note" in combined.lower(), "Notes text should be present in content"
        finally:
            os.unlink(tmp_path)

    # --- Empty slide skipping ---

    def test_empty_slides_skipped(self):
        path = _PPTX_DIR / "empty_slides.pptx"
        if not path.exists():
            pytest.skip("empty_slides.pptx not found")
        chunks = _get(path, "default")
        # Should have ≤4 chunks (title + 3 content; 2 blank skipped)
        assert len(chunks) <= 4, \
            f"Expected ≤4 chunks (2 blank slides should be skipped), got {len(chunks)}"

    # --- Prose heavy: semantic merges slides ---

    def test_semantic_merges_related_slides(self):
        path = _PPTX_DIR / "prose_heavy.pptx"
        if not path.exists():
            pytest.skip("prose_heavy.pptx not found")
        structural_n = len(_get(path, "default"))
        semantic_n   = len(_get(path, "semantic"))
        assert semantic_n <= structural_n, \
            "Semantic mode should produce ≤ structural chunk count for prose deck"

    # --- Section mode groups content ---

    def test_section_mode_produces_fewer_chunks_than_default(self):
        path = _PPTX_DIR / "sectioned.pptx"
        if not path.exists():
            pytest.skip("sectioned.pptx not found")
        default_n = len(_get(path, "default"))
        # Section mode may emit heading chunks + body — but body count < default
        body_chunks = [c for c in _get(path, "section") if c["content_type"] != "heading"]
        assert len(body_chunks) < default_n, \
            "Section body chunks should be fewer than individual slide chunks"

    # --- Sliding window covers all slides ---

    def test_sliding_window_covers_all_content(self):
        path = _PPTX_DIR / "minimal.pptx"
        if not path.exists():
            pytest.skip("minimal.pptx not found")
        chunks = _get(path, "sliding_window", window_size=2, overlap=0)
        all_content = " ".join(c["content"] for c in chunks)
        # At minimum the intro and conclusion words should appear somewhere
        assert "Introduction" in all_content or "introduction" in all_content.lower()

    # --- Sentence mode: every sentence ends with punctuation or is non-empty ---

    def test_sentence_chunks_non_empty(self):
        path = _PPTX_DIR / "prose_heavy.pptx"
        if not path.exists():
            pytest.skip("prose_heavy.pptx not found")
        for c in _get(path, "sentence", sentences_per_chunk=2):
            assert c["content"].strip(), "Sentence chunk must not be empty"

    # --- Average chunk lengths by mode (quality signal) ---

    @pytest.mark.parametrize("path", [
        p for p in ALL_PPTX if p.name == "large_deck.pptx"
    ], ids=["large_deck.pptx"])
    def test_avg_chunk_length_reasonable(self, path):
        for mode in ("default", "semantic", "sentence", "page_aware"):
            kw = {"sentences_per_chunk": 3} if mode == "sentence" else {}
            kw = {"paragraphs_per_page": 5} if mode == "page_aware" else kw
            chunks = _get(path, mode, **kw)
            avg = _avg_chunk_len(chunks)
            # Chunks should be of meaningful size — not trivially tiny
            assert avg >= 30, f"Average chunk length too small ({avg:.0f} chars) for mode={mode}"


# ═══════════════════════════════════════════════════════════════════════════════
# 5. SLIDING WINDOW — overlap / step correctness
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlidingWindowCorrectness:

    def _get_sw(self, path, ws, ov):
        return _get(path, "sliding_window", window_size=ws, overlap=ov)

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_window_size_1_no_overlap_equals_structural(self, path):
        sw   = self._get_sw(path, 1, 0)
        st   = _get(path, "default")
        # Each window should correspond 1:1 with structural (same slide count)
        assert len(sw) == len(st), (
            f"window_size=1,overlap=0 should give same chunk count as structural:"
            f" {len(sw)} vs {len(st)}"
        )

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_larger_window_fewer_chunks(self, path):
        sw2 = self._get_sw(path, 2, 0)
        sw3 = self._get_sw(path, 3, 0)
        if len(_get(path, "default")) >= 3:
            assert len(sw3) <= len(sw2), "Larger window should produce ≤ chunk count"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_overlap_increases_chunk_count(self, path):
        sw_no_ov = self._get_sw(path, 3, 0)
        sw_ov    = self._get_sw(path, 3, 1)
        n = len(_get(path, "default"))
        if n >= 4:
            assert len(sw_ov) >= len(sw_no_ov), \
                "overlap=1 should give ≥ chunk count vs overlap=0"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_slide_ranges_advance_monotonically(self, path):
        for c1, c2 in zip(
            _get(path, "sliding_window", window_size=3, overlap=1),
            _get(path, "sliding_window", window_size=3, overlap=1)[1:],
        ):
            r1 = c1["metadata"]["slide_range"][0]
            r2 = c2["metadata"]["slide_range"][0]
            assert r2 > r1, f"slide_range start must increase: {r1} → {r2}"


# ═══════════════════════════════════════════════════════════════════════════════
# 6. DEFAULT vs STRUCTURAL EQUIVALENCE
# ═══════════════════════════════════════════════════════════════════════════════

class TestDefaultStructuralEquivalence:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_same_chunk_count(self, path):
        d = _get(path, "default")
        s = _get(path, "structural")
        assert len(d) == len(s)

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_same_content(self, path):
        d = _get(path, "default")
        s = _get(path, "structural")
        for dc, sc in zip(d, s):
            assert dc["content"] == sc["content"]


# ═══════════════════════════════════════════════════════════════════════════════
# 7. STREAMING — PptxStreamIterator
# ═══════════════════════════════════════════════════════════════════════════════

class TestStreaming:
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_stream_default_matches_batch(self, path):
        batch   = _get(path, "default")
        streamed = list(stream_chunks(str(path), mode="default"))
        assert len(streamed) == len(batch)
        for bc, sc in zip(batch, streamed):
            assert bc["content"] == sc["content"]

    @pytest.mark.parametrize("mode", [
        "default", "structural", "semantic", "section",
        "sliding_window", "sentence", "page_aware",
    ])
    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_stream_all_modes_schema(self, path, mode):
        kw = {}
        if mode == "sliding_window":
            kw = {"window_size": 3, "overlap": 1}
        elif mode == "sentence":
            kw = {"sentences_per_chunk": 3}
        elif mode == "page_aware":
            kw = {"paragraphs_per_page": 3}
        streamed = list(stream_chunks(str(path), mode=mode, **kw))
        _assert_schema(streamed, f"stream/{mode}/{path.name}")

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_stream_yields_one_at_a_time(self, path):
        it = stream_chunks(str(path), mode="default")
        first = next(it)
        assert isinstance(first, dict)
        _assert_schema([first], "stream_first_chunk")


# ═══════════════════════════════════════════════════════════════════════════════
# 8. ERROR HANDLING / VALIDATION
# ═══════════════════════════════════════════════════════════════════════════════

class TestErrorHandling:

    def test_missing_file_raises(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            chunk_pptx(str(tmp_path / "nonexistent.pptx"))

    def test_wrong_extension_raises(self, tmp_path):
        f = tmp_path / "doc.docx"
        f.write_bytes(b"fake")
        with pytest.raises(ValueError, match="pptx"):
            chunk_pptx(str(f))

    def test_bad_mode_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError, match="mode"):
            chunk_pptx(str(path), mode="nonexistent_mode")

    def test_invalid_window_size_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            chunk_pptx(str(path), mode="sliding_window", window_size=0)

    def test_overlap_gte_window_size_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            chunk_pptx(str(path), mode="sliding_window", window_size=2, overlap=2)

    def test_sentences_per_chunk_zero_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            chunk_pptx(str(path), mode="sentence", sentences_per_chunk=0)

    def test_paragraphs_per_page_zero_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            chunk_pptx(str(path), mode="page_aware", paragraphs_per_page=0)

    def test_corrupted_pptx_raises(self, tmp_path):
        bad = tmp_path / "bad.pptx"
        bad.write_bytes(b"not a zip at all")
        with pytest.raises(Exception):
            chunk_pptx(str(bad))


# ═══════════════════════════════════════════════════════════════════════════════
# 9. PERFORMANCE BENCHMARKS
# ═══════════════════════════════════════════════════════════════════════════════

class TestPerformance:
    """Timing thresholds: Rust path should be fast; Python overhead should be minimal."""

    MODES_WITH_PARAMS = [
        ("default", {}),
        ("semantic", {}),
        ("section", {}),
        ("sliding_window", {"window_size": 3, "overlap": 1}),
        ("sentence", {"sentences_per_chunk": 3}),
        ("page_aware", {"paragraphs_per_page": 5}),
    ]

    @pytest.mark.parametrize("mode,kw", MODES_WITH_PARAMS,
                              ids=[m for m, _ in MODES_WITH_PARAMS])
    def test_large_deck_rust_under_200ms(self, mode, kw):
        path = _PPTX_DIR / "large_deck.pptx"
        if not path.exists():
            pytest.skip("large_deck.pptx not found")
        _, rust_ms, _, _ = _timed_get(path, mode, **kw)
        assert rust_ms < 200, f"Rust time too slow: {rust_ms:.1f}ms for mode={mode}"

    @pytest.mark.parametrize("mode,kw", MODES_WITH_PARAMS,
                              ids=[m for m, _ in MODES_WITH_PARAMS])
    def test_large_deck_wall_under_500ms(self, mode, kw):
        path = _PPTX_DIR / "large_deck.pptx"
        if not path.exists():
            pytest.skip("large_deck.pptx not found")
        _, _, _, wall_ms = _timed_get(path, mode, **kw)
        assert wall_ms < 500, f"Wall time too slow: {wall_ms:.1f}ms for mode={mode}"

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_timing_keys_present(self, path):
        _, timing = chunk_pptx(str(path))
        assert "rust_ms" in timing
        assert "python_ms" in timing
        assert isinstance(timing["rust_ms"], (int, float))
        assert isinstance(timing["python_ms"], (int, float))

    @pytest.mark.parametrize("path", ALL_PPTX, ids=PPTX_IDS)
    def test_rust_ms_positive(self, path):
        _, timing = chunk_pptx(str(path))
        assert timing["rust_ms"] > 0, "rust_ms must be > 0"


class TestPptxSourceEquivalence:
    def test_bytes_equals_path(self):
        path = ALL_PPTX[0]
        by_path = get_chunks(str(path), mode="structural")
        by_bytes = get_chunks_from_bytes(
            Path(path).read_bytes(), filename="x.pptx", mode="structural"
        )
        assert [c["content"] for c in by_bytes] == [c["content"] for c in by_path]

    def test_fileobj_equals_path(self):
        path = ALL_PPTX[0]
        by_path = get_chunks(str(path), mode="structural")
        with open(path, "rb") as f:
            by_fileobj = get_chunks_from_fileobj(f, filename="x.pptx", mode="structural")
        assert [c["content"] for c in by_fileobj] == [c["content"] for c in by_path]

    def test_upload_equals_path(self):
        path = ALL_PPTX[0]
        by_path = get_chunks(str(path), mode="structural")
        upload = _DummyUpload(Path(path).read_bytes(), "x.pptx")
        by_upload = get_chunks_from_upload(upload, mode="structural")
        assert [c["content"] for c in by_upload] == [c["content"] for c in by_path]


class TestPptxStreamingValidation:
    def test_sliding_window_window_size_zero_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            stream_chunk_pptx(str(path), mode="sliding_window", window_size=0)

    def test_sliding_window_overlap_gte_window_size_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            stream_chunk_pptx(str(path), mode="sliding_window", window_size=3, overlap=3)

    def test_sentence_sentences_per_chunk_zero_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises(ValueError):
            stream_chunk_pptx(str(path), mode="sentence", sentences_per_chunk=0)

    def test_invalid_mode_raises(self):
        path = ALL_PPTX[0]
        with pytest.raises((ValueError, NotImplementedError)):
            stream_chunk_pptx(str(path), mode="not_a_real_mode")
